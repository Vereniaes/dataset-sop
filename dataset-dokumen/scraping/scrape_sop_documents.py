"""
scrape_sop_documents.py
=======================
Download dan ekstrak teks dari dokumen SOP (PDF, DOCX, PPT/PPTX)
lalu simpan ke folder terpisah berdasarkan tipe file.

Struktur output:
  file-documents/
    pdf/
      fnb/                → PDF kategori F&B
      retail/             → PDF kategori Retail
      .../{category}/     → per kategori
    docx/
      .../{category}/
    ppt/
      .../{category}/
    extracted/
      fnb.jsonl           → hasil ekstrak per kategori
      retail.jsonl
      all.jsonl           → gabungan semua kategori

Cara pakai:
  1. Isi URL dokumen di file-documents/url_list.jsonl
  2. Jalankan: python scrape_sop_documents.py
  3. (Opsional) jalankan dengan --normalize untuk konversi ke format 7 seksi

Dependensi:
  pip install pdfplumber python-docx python-pptx requests tqdm
  (opsional OCR): pip install pytesseract pillow pdf2image
"""

import os
import json
import time
import hashlib
import logging
import argparse
import requests
from pathlib import Path
from urllib.parse import urlparse
from typing import Optional

# ─── Konfigurasi ──────────────────────────────────────────────────────────────

BASE_DIR      = Path(__file__).parent / "file-documents"
PDF_DIR       = BASE_DIR / "pdf"
DOCX_DIR      = BASE_DIR / "docx"
PPT_DIR       = BASE_DIR / "ppt"
EXTRACTED_DIR = BASE_DIR / "extracted"
URL_LIST      = BASE_DIR / "url_list.jsonl"

# Root folder dibuat di sini; subfolder category dibuat dinamis saat runtime
for d in [PDF_DIR, DOCX_DIR, PPT_DIR, EXTRACTED_DIR]:
    d.mkdir(parents=True, exist_ok=True)

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s",
    handlers=[
        logging.FileHandler(BASE_DIR / "scrape_log.txt", encoding="utf-8"),
        logging.StreamHandler()
    ]
)
log = logging.getLogger(__name__)

# ─── Helper: Download ──────────────────────────────────────────────────────────

HEADERS = {
    "User-Agent": (
        "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
        "AppleWebKit/537.36 (KHTML, like Gecko) "
        "Chrome/124.0.0.0 Safari/537.36"
    )
}


def get_dest_dir(file_type: str, category: str) -> Path:
    """Return folder tujuan: {filetype}/{category}/ dan buat jika belum ada."""
    root = {"pdf": PDF_DIR, "docx": DOCX_DIR, "ppt": PPT_DIR}.get(file_type, PDF_DIR)
    dest = root / category
    dest.mkdir(parents=True, exist_ok=True)
    return dest


def download_file(url: str, dest_dir: Path, filename: str) -> Optional[Path]:
    """Download file dari URL ke folder tujuan. Return path jika berhasil."""
    dest_path = dest_dir / filename
    if dest_path.exists():
        log.info(f"  ✓ Sudah ada: {filename}, skip download")
        return dest_path

    try:
        log.info(f"  ↓ Download: {url}")
        resp = requests.get(url, headers=HEADERS, timeout=30, stream=True)
        resp.raise_for_status()

        with open(dest_path, "wb") as f:
            for chunk in resp.iter_content(chunk_size=8192):
                f.write(chunk)

        size_kb = dest_path.stat().st_size / 1024
        log.info(f"  ✓ Tersimpan: {filename} ({size_kb:.1f} KB)")
        return dest_path

    except requests.RequestException as e:
        log.error(f"  ✗ Gagal download {url}: {e}")
        return None


def get_safe_filename(url: str, category: str, ext: str) -> str:
    """Buat nama file yang aman dari URL + kategori."""
    url_hash = hashlib.md5(url.encode()).hexdigest()[:8]
    return f"{category}_{url_hash}{ext}"


# ─── Ekstraksi PDF ─────────────────────────────────────────────────────────────

def extract_pdf(path: Path) -> str:
    """Ekstrak teks dari PDF menggunakan pdfplumber.
    Fallback ke pytesseract jika PDF adalah scan (teks kosong)."""
    try:
        import pdfplumber
    except ImportError:
        log.error("pdfplumber belum terinstall. Jalankan: pip install pdfplumber")
        return ""

    text_parts = []
    try:
        with pdfplumber.open(path) as pdf:
            for i, page in enumerate(pdf.pages):
                page_text = page.extract_text()
                if page_text and page_text.strip():
                    text_parts.append(page_text.strip())
                else:
                    log.warning(f"  ⚠ Halaman {i+1} kosong (mungkin scan) di {path.name}")

        if not text_parts:
            log.warning(f"  ⚠ Tidak ada teks terdeteksi di {path.name}, coba OCR...")
            return _extract_pdf_ocr(path)

        return "\n\n".join(text_parts)

    except Exception as e:
        log.error(f"  ✗ Error ekstrak PDF {path.name}: {e}")
        return ""


def _extract_pdf_ocr(path: Path) -> str:
    """Fallback OCR untuk PDF scan (butuh pytesseract + poppler)."""
    try:
        from pdf2image import convert_from_path
        import pytesseract
    except ImportError:
        log.warning("  OCR library tidak tersedia (pip install pytesseract pdf2image)")
        return ""

    try:
        images = convert_from_path(str(path), dpi=200)
        texts = [pytesseract.image_to_string(img, lang="ind+eng") for img in images]
        return "\n\n".join(t.strip() for t in texts if t.strip())
    except Exception as e:
        log.error(f"  ✗ Error OCR {path.name}: {e}")
        return ""


# ─── Ekstraksi DOCX ────────────────────────────────────────────────────────────

def extract_docx(path: Path) -> str:
    """Ekstrak teks dari DOCX menggunakan python-docx."""
    try:
        from docx import Document
    except ImportError:
        log.error("python-docx belum terinstall. Jalankan: pip install python-docx")
        return ""

    try:
        doc = Document(str(path))
        parts = []

        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text.strip())

        # Ekstrak teks dari tabel juga
        for table in doc.tables:
            for row in table.rows:
                row_texts = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if row_texts:
                    parts.append(" | ".join(row_texts))

        return "\n".join(parts)

    except Exception as e:
        log.error(f"  ✗ Error ekstrak DOCX {path.name}: {e}")
        return ""


# ─── Ekstraksi PPT/PPTX ────────────────────────────────────────────────────────

def extract_ppt(path: Path) -> str:
    """Ekstrak teks dari PPT/PPTX menggunakan python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        log.error("python-pptx belum terinstall. Jalankan: pip install python-pptx")
        return ""

    try:
        prs = Presentation(str(path))
        parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)

            if slide_texts:
                parts.append(f"[Slide {slide_num}]\n" + "\n".join(slide_texts))

        return "\n\n".join(parts)

    except Exception as e:
        log.error(f"  ✗ Error ekstrak PPT {path.name}: {e}")
        return ""


# ─── Router Ekstraksi ──────────────────────────────────────────────────────────

EXTRACTORS = {
    "pdf":  extract_pdf,
    "docx": extract_docx,
    "ppt":  extract_ppt,
}

# TARGET_DIRS tidak lagi statis — pakai get_dest_dir(file_type, category)

EXT_MAP = {
    "pdf":  ".pdf",
    "docx": ".docx",
    "ppt":  ".pptx",
}


def is_valid_text(text: str, min_chars: int = 200) -> bool:
    """Cek apakah teks hasil ekstraksi layak (tidak kosong/terlalu pendek)."""
    clean = text.strip()
    return len(clean) >= min_chars


# ─── Pipeline Utama ────────────────────────────────────────────────────────────

def load_url_list() -> list[dict]:
    """Muat daftar URL dari url_list.jsonl."""
    if not URL_LIST.exists():
        log.error(f"File URL tidak ditemukan: {URL_LIST}")
        return []

    entries = []
    with open(URL_LIST, "r", encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if not line or line.startswith("//"):
                continue
            try:
                entry = json.loads(line)
                # Skip template placeholder
                if "example.com" in entry.get("url", ""):
                    log.warning(f"  ⚠ Skip placeholder URL: {entry['url']}")
                    continue
                entries.append(entry)
            except json.JSONDecodeError as e:
                log.warning(f"  ⚠ Skip baris tidak valid: {e}")

    log.info(f"Loaded {len(entries)} URL dari {URL_LIST.name}")
    return entries


def process_entry(entry: dict) -> Optional[dict]:
    """Proses satu entry: download → ekstrak → return dict hasil."""
    url       = entry.get("url", "").strip()
    category  = entry.get("category", "unknown")
    label     = entry.get("label", category)
    file_type = entry.get("file_type", "pdf").lower()

    if file_type not in EXT_MAP:
        log.warning(f"  ⚠ Tipe file tidak dikenali: {file_type} untuk {url}")
        return None

    # 1. Download ke {filetype}/{category}/
    ext       = EXT_MAP[file_type]
    filename  = get_safe_filename(url, category, ext)
    dest_dir  = get_dest_dir(file_type, category)   # pdf/fnb/, docx/retail/, dst
    file_path = download_file(url, dest_dir, filename)

    if not file_path:
        return None

    time.sleep(1)  # delay agar tidak kena rate limit

    # 2. Ekstrak teks
    extractor = EXTRACTORS[file_type]
    text = extractor(file_path)

    if not is_valid_text(text):
        log.warning(f"  ⚠ Teks terlalu pendek atau kosong: {filename} ({len(text)} chars)")
        return None

    log.info(f"  ✓ Ekstrak OK: {filename} ({len(text)} chars)")

    return {
        "category":   category,
        "label":      label,
        "source_url": url,
        "file_type":  file_type,
        "filename":   filename,
        "char_count": len(text),
        "text":       text,
    }


def run(dry_run: bool = False):
    """Jalankan pipeline scraping dokumen SOP."""
    entries = load_url_list()
    if not entries:
        log.error("Tidak ada URL valid ditemukan. Isi url_list.jsonl terlebih dahulu.")
        return

    results = []
    skipped = 0
    failed  = 0

    for i, entry in enumerate(entries, 1):
        url = entry.get("url", "")
        cat = entry.get("category", "?")
        log.info(f"\n[{i}/{len(entries)}] Kategori: {cat}")
        log.info(f"  URL: {url}")

        if dry_run:
            log.info("  [DRY RUN] Skip proses")
            continue

        result = process_entry(entry)

        if result:
            results.append(result)
        else:
            failed += 1

    # 3. Simpan hasil extracted per kategori + all.jsonl
    if results:
        from collections import defaultdict, Counter

        # Per kategori: extracted/{category}.jsonl
        per_category: dict = defaultdict(list)
        for r in results:
            per_category[r["category"]].append(r)

        for cat, items in per_category.items():
            cat_path = EXTRACTED_DIR / f"{cat}.jsonl"
            with open(cat_path, "w", encoding="utf-8") as f:
                for item in items:
                    f.write(json.dumps(item, ensure_ascii=False) + "\n")
            log.info(f"  ✓ Saved: extracted/{cat}.jsonl ({len(items)} dokumen)")

        # Gabungan: extracted/all.jsonl
        all_path = EXTRACTED_DIR / "all.jsonl"
        with open(all_path, "w", encoding="utf-8") as f:
            for r in results:
                f.write(json.dumps(r, ensure_ascii=False) + "\n")

        log.info(f"\n{'='*60}")
        log.info(f"✅ Selesai!")
        log.info(f"  Berhasil   : {len(results)} dokumen")
        log.info(f"  Gagal      : {failed} dokumen")
        log.info(f"  Output all : extracted/all.jsonl")
        log.info(f"  PDF folder : pdf/{{category}}/")
        log.info(f"  DOCX folder: docx/{{category}}/")
        log.info(f"  PPT folder : ppt/{{category}}/")

        log.info("\n📊 Ringkasan per Kategori:")
        cats = Counter(r["category"] for r in results)
        for cat, count in sorted(cats.items()):
            log.info(f"  [{cat}] {count} dokumen → extracted/{cat}.jsonl")
    else:
        log.warning("\n⚠ Tidak ada dokumen yang berhasil diekstrak.")


# ─── CLI ───────────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Download & ekstrak dokumen SOP (PDF/DOCX/PPT)"
    )
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="Test tanpa download/ekstrak (hanya validasi URL list)"
    )
    parser.add_argument(
        "--category",
        type=str,
        default=None,
        help="Filter hanya proses kategori tertentu (contoh: fnb)"
    )
    args = parser.parse_args()

    # Filter kategori jika ada
    if args.category:
        original_load = load_url_list
        def load_url_list():  # noqa: F811
            entries = original_load()
            filtered = [e for e in entries if e.get("category") == args.category]
            log.info(f"Filter kategori '{args.category}': {len(filtered)} URL")
            return filtered

    run(dry_run=args.dry_run)
